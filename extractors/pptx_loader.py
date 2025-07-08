from pptx import Presentation
from io import BytesIO

def load_pptx_text(file_bytes: bytes) -> str:
    prs = Presentation(BytesIO(file_bytes))
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)